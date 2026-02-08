"""
Document Processing Task V2 - Using Microservices
OCR: GLM-OCR Service (Port 8001)
Chunking: Hybrid Chunking Service (Port 8003)
Embedding: E5 Embedding Service (Port 8002)
"""
import os
from pathlib import Path
from celery import Task
from celery_app import app
from typing import List, Dict, Any
import logging
import httpx
import asyncio

logger = logging.getLogger(__name__)

# Service URLs
OCR_SERVICE_URL = os.getenv("OCR_SERVICE_URL", "http://ocr_service:8001")
EMBEDDING_SERVICE_URL = os.getenv("EMBEDDING_SERVICE_URL", "http://embedding_service:8002")
CHUNKING_SERVICE_URL = os.getenv("CHUNKING_SERVICE_URL", "http://chunking_service:8003")


async def extract_text_async(file_path: str, file_type: str) -> str:
    """Extract text from document using appropriate service"""
    text = ""
    
    try:
        if file_type in [".tif", ".tiff", ".png", ".jpg", ".jpeg"]:
            # Use GLM-OCR Service for images
            async with httpx.AsyncClient(timeout=120.0) as client:
                with open(file_path, "rb") as f:
                    files = {"file": (Path(file_path).name, f, "image/png")}
                    data = {"language": "en"}  # Can be dynamic: en, ko, zh, ja
                    
                    response = await client.post(
                        f"{OCR_SERVICE_URL}/ocr",
                        files=files,
                        data=data
                    )
                    response.raise_for_status()
                    result = response.json()
                    text = result.get("text", "")
                    logger.info(f"✅ OCR extracted {len(text)} chars from {file_path}")
        
        elif file_type in [".pdf"]:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            text = "\n".join([page.extract_text() for page in reader.pages])
        
        elif file_type in [".docx", ".doc"]:
            import docx
            doc = docx.Document(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
        
        elif file_type in [".xlsx", ".xls"]:
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            text = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    text += " ".join([str(cell) for cell in row if cell]) + "\n"
        
        elif file_type in [".pptx", ".ppt"]:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        
        else:
            logger.warning(f"Unsupported file type: {file_type}")
    
    except Exception as e:
        logger.error(f"Failed to extract text from {file_path}: {e}")
    
    return text.strip()


async def chunk_text_async(
    text: str,
    method: str = "hybrid",
    chunk_size: int = 1000,
    overlap: int = 200
) -> List[str]:
    """Chunk text using Hybrid Chunking Service"""
    if not text:
        return []
    
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            response = await client.post(
                f"{CHUNKING_SERVICE_URL}/chunk",
                json={
                    "text": text,
                    "method": method,
                    "chunk_size": chunk_size,
                    "chunk_overlap": overlap
                }
            )
            response.raise_for_status()
            result = response.json()
            chunks = result.get("chunks", [])
            logger.info(f"✅ Chunking created {len(chunks)} chunks (method={method})")
            return chunks
    
    except Exception as e:
        logger.error(f"Chunking failed: {e}")
        return []


async def embed_chunks_async(chunks: List[str], batch_size: int = 32) -> List[List[float]]:
    """Generate embeddings using E5 Embedding Service"""
    if not chunks:
        return []
    
    try:
        async with httpx.AsyncClient(timeout=180.0) as client:
            response = await client.post(
                f"{EMBEDDING_SERVICE_URL}/embed",
                json={
                    "texts": chunks,
                    "normalize": True,
                    "batch_size": batch_size
                }
            )
            response.raise_for_status()
            result = response.json()
            embeddings = result.get("embeddings", [])
            dimension = result.get("dimension", 0)
            logger.info(f"✅ Embedded {len(embeddings)} chunks (dimension={dimension})")
            return embeddings
    
    except Exception as e:
        logger.error(f"Embedding failed: {e}")
        return []


@app.task(name="tasks.document_processing_v2.process_document", bind=True)
def process_document(
    self: Task,
    file_path: str,
    file_hash: str,
    department: str,
    role: str
):
    """
    Process a single document using microservices:
    1. Extract text (GLM-OCR for images, libraries for docs)
    2. Chunk text (Hybrid Chunking Service)
    3. Generate embeddings (E5 Embedding Service)
    4. Upsert to Qdrant
    5. Log to PostgreSQL
    """
    logger.info(f"Processing document: {file_path}")
    
    try:
        # Get file info
        path = Path(file_path)
        filename = path.name
        file_type = path.suffix.lower()
        file_size = path.stat().st_size
        
        # Step 1: Extract text (async)
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        text = loop.run_until_complete(extract_text_async(file_path, file_type))
        
        if not text:
            logger.warning(f"No text extracted from {file_path}")
            return {"status": "skipped", "reason": "no_text"}
        
        # Step 2: Chunk text (async)
        chunks = loop.run_until_complete(
            chunk_text_async(
                text,
                method="hybrid",
                chunk_size=1000,
                overlap=200
            )
        )
        
        if not chunks:
            logger.warning(f"No chunks created from {file_path}")
            return {"status": "skipped", "reason": "no_chunks"}
        
        logger.info(f"Generated {len(chunks)} chunks from {filename}")
        
        # Step 3: Generate embeddings (async)
        embeddings = loop.run_until_complete(embed_chunks_async(chunks, batch_size=32))
        
        if not embeddings:
            logger.warning(f"No embeddings generated for {file_path}")
            return {"status": "skipped", "reason": "no_embeddings"}
        
        # Step 4: Upsert to Qdrant
        # TODO: Import and use qdrant_service
        # from backend.app.services.qdrant_service import qdrant_service
        # point_ids = qdrant_service.upsert_documents(...)
        
        logger.info(f"Successfully processed {filename}")
        
        loop.close()
        
        return {
            "status": "completed",
            "filename": filename,
            "chunks": len(chunks),
            "embeddings": len(embeddings),
            "department": department,
            "role": role
        }
    
    except Exception as e:
        logger.error(f"Failed to process document {file_path}: {e}")
        return {
            "status": "failed",
            "error": str(e)
        }
