"""
Document Processing Task - Using Microservices

OCR: GLM-OCR Service (Port 8001)
Chunking: Hybrid Chunking Service (Port 8003)
Embedding: E5 Embedding Service (Port 8002)
"""
import asyncio
import logging
import os
from pathlib import Path
from typing import Any, Dict, List, Optional
from uuid import uuid4

from celery import Task
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, PointIdsList, PointStruct, VectorParams
from sqlalchemy import create_engine, text

from celery_app import app

from shared.service_client import OCRClient, ChunkingClient, EmbeddingClient

logger = logging.getLogger(__name__)

# Service clients
ocr_client = OCRClient(base_url=os.getenv("OCR_URL", "http://ocr_service:8001"))
chunking_client = ChunkingClient(base_url=os.getenv("CHUNKING_URL", "http://chunking_service:8003"))
embedding_client = EmbeddingClient(base_url=os.getenv("EMBEDDING_URL", "http://embedding_service:8002"))

_db_engine = None


def get_db_engine():
    """Create and cache SQLAlchemy engine for worker metadata writes."""
    global _db_engine
    if _db_engine is None:
        host = os.getenv("POSTGRES_HOST", "postgres")
        port = os.getenv("POSTGRES_PORT", "5432")
        user = os.getenv("POSTGRES_USER", "admin")
        password = os.getenv("POSTGRES_PASSWORD", "securepassword")
        database = os.getenv("POSTGRES_DB", "onprem_llm")
        dsn = f"postgresql+psycopg2://{user}:{password}@{host}:{port}/{database}"
        _db_engine = create_engine(dsn, pool_pre_ping=True)
    return _db_engine


def get_qdrant_client() -> QdrantClient:
    """Create Qdrant client from environment."""
    return QdrantClient(
        host=os.getenv("QDRANT_HOST", "qdrant"),
        port=int(os.getenv("QDRANT_PORT", "6333")),
        timeout=30,
    )


def ensure_qdrant_collection(client: QdrantClient) -> str:
    """Create collection if missing."""
    collection_name = os.getenv("QDRANT_COLLECTION_NAME", "documents")
    vector_size = int(os.getenv("QDRANT_VECTOR_SIZE", "1024"))
    collections = client.get_collections().collections
    names = {collection.name for collection in collections}
    if collection_name not in names:
        client.create_collection(
            collection_name=collection_name,
            vectors_config=VectorParams(size=vector_size, distance=Distance.COSINE),
        )
    return collection_name


def get_existing_document(file_hash: str) -> Optional[Dict[str, Any]]:
    """Load existing document row by file hash."""
    engine = get_db_engine()
    with engine.connect() as conn:
        row = conn.execute(
            text(
                """
                SELECT doc_id
                FROM document
                WHERE hash = :file_hash
                """
            ),
            {"file_hash": file_hash},
        ).mappings().first()
        return dict(row) if row else None


def get_existing_qdrant_ids(doc_id: int) -> List[str]:
    """Get existing Qdrant point IDs for a document from doc_chunk table."""
    engine = get_db_engine()
    with engine.connect() as conn:
        rows = conn.execute(
            text("SELECT qdrant_id FROM doc_chunk WHERE doc_id = :doc_id AND qdrant_id IS NOT NULL"),
            {"doc_id": doc_id},
        ).all()
        return [row[0] for row in rows]


def delete_existing_qdrant_points(client: QdrantClient, collection_name: str, point_ids: List[str]) -> None:
    """Delete previous vectors before re-indexing modified document."""
    if not point_ids:
        return
    client.delete(
        collection_name=collection_name,
        points_selector=PointIdsList(points=point_ids),
    )


def delete_existing_chunks(doc_id: int) -> None:
    """Delete existing chunks from doc_chunk table before re-indexing."""
    engine = get_db_engine()
    with engine.begin() as conn:
        conn.execute(
            text("DELETE FROM doc_chunk WHERE doc_id = :doc_id"),
            {"doc_id": doc_id},
        )


def upsert_to_qdrant(
    client: QdrantClient,
    collection_name: str,
    doc_id: int,
    chunks: List[str],
    embeddings: List[List[float]],
    dept_id: int,
    role_id: int,
    filename: str,
    file_path: str,
    file_type: str,
) -> List[str]:
    """Upsert chunk vectors with RBAC payload into Qdrant."""
    points: List[PointStruct] = []
    point_ids: List[str] = []
    for idx, (chunk, embedding) in enumerate(zip(chunks, embeddings)):
        point_id = str(uuid4())
        point_ids.append(point_id)
        points.append(
            PointStruct(
                id=point_id,
                vector=embedding,
                payload={
                    "document_id": doc_id,
                    "chunk_index": idx,
                    "content": chunk,
                    "filename": filename,
                    "file_path": file_path,
                    "file_type": file_type,
                    "dept_id": dept_id,
                    "role_id": role_id,
                },
            )
        )

    client.upsert(collection_name=collection_name, points=points)
    return point_ids


def upsert_document_metadata(
    filename: str,
    file_path: str,
    file_type: str,
    file_size: int,
    file_hash: str,
    dept_id: int,
    role_id: int,
    existing_doc_id: Optional[int] = None,
) -> int:
    """Upsert document metadata into PostgreSQL. Returns doc_id."""
    engine = get_db_engine()
    with engine.begin() as conn:
        if existing_doc_id:
            conn.execute(
                text(
                    """
                    UPDATE document SET
                        file_name = :file_name,
                        path = :path,
                        type = :type,
                        size = :size,
                        status = 'processing',
                        updated_at = NOW()
                    WHERE doc_id = :doc_id
                    """
                ),
                {
                    "file_name": filename,
                    "path": file_path,
                    "type": file_type,
                    "size": file_size,
                    "doc_id": existing_doc_id,
                },
            )
            return existing_doc_id
        else:
            result = conn.execute(
                text(
                    """
                    INSERT INTO document (
                        file_name, path, type, hash, size,
                        dept_id, role_id, status, created_at, updated_at
                    )
                    VALUES (
                        :file_name, :path, :type, :hash, :size,
                        :dept_id, :role_id, 'processing', NOW(), NOW()
                    )
                    RETURNING doc_id
                    """
                ),
                {
                    "file_name": filename,
                    "path": file_path,
                    "type": file_type,
                    "hash": file_hash,
                    "size": file_size,
                    "dept_id": dept_id,
                    "role_id": role_id,
                },
            )
            return result.scalar_one()


def insert_doc_chunks(
    doc_id: int,
    chunks: List[str],
    point_ids: List[str],
    embed_model: str = "e5-large",
) -> None:
    """Insert chunk rows into doc_chunk table."""
    engine = get_db_engine()
    with engine.begin() as conn:
        for idx, (chunk, point_id) in enumerate(zip(chunks, point_ids)):
            conn.execute(
                text(
                    """
                    INSERT INTO doc_chunk (
                        doc_id, chunk_idx, content, token_cnt, qdrant_id, embed_model, created_at, updated_at
                    )
                    VALUES (
                        :doc_id, :chunk_idx, :content, :token_cnt, :qdrant_id, :embed_model, NOW(), NOW()
                    )
                    """
                ),
                {
                    "doc_id": doc_id,
                    "chunk_idx": idx,
                    "content": chunk,
                    "token_cnt": len(chunk.split()),
                    "qdrant_id": point_id,
                    "embed_model": embed_model,
                },
            )


def mark_document_status(doc_id: int, status: str) -> None:
    """Update document status (indexed / failed)."""
    engine = get_db_engine()
    with engine.begin() as conn:
        conn.execute(
            text("UPDATE document SET status = :status, updated_at = NOW() WHERE doc_id = :doc_id"),
            {"status": status, "doc_id": doc_id},
        )


async def extract_text_async(file_path: str, file_type: str) -> str:
    """Extract text from document using appropriate service."""
    extracted = ""

    try:
        if file_type in [".tif", ".tiff", ".png", ".jpg", ".jpeg"]:
            result = await ocr_client.ocr(file_path)
            extracted = result.get("text", "")
            logger.info(f"OCR extracted {len(extracted)} chars from {file_path}")

        elif file_type in [".pdf"]:
            from PyPDF2 import PdfReader
            reader = PdfReader(file_path)
            extracted = "\n".join([page.extract_text() for page in reader.pages])

        elif file_type in [".docx", ".doc"]:
            import docx
            doc = docx.Document(file_path)
            extracted = "\n".join([para.text for para in doc.paragraphs])

        elif file_type in [".xlsx", ".xls"]:
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            extracted = ""
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    extracted += " ".join([str(cell) for cell in row if cell]) + "\n"

        elif file_type in [".pptx", ".ppt"]:
            from pptx import Presentation
            prs = Presentation(file_path)
            extracted = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        extracted += shape.text + "\n"

        else:
            logger.warning(f"Unsupported file type: {file_type}")

    except Exception as e:
        logger.error(f"Failed to extract text from {file_path}: {e}")

    return extracted.strip()


async def chunk_text_async(
    text: str,
    method: str = "hybrid",
    chunk_size: int = 1000,
    overlap: int = 200
) -> List[str]:
    """Chunk text using Hybrid Chunking Service."""
    if not text:
        return []

    try:
        result = await chunking_client.chunk(
            text=text,
            method=method,
            chunk_size=chunk_size,
            chunk_overlap=overlap,
        )
        chunks = result.get("chunks", [])
        logger.info(f"Chunking created {len(chunks)} chunks (method={method})")
        return chunks

    except Exception as e:
        logger.error(f"Chunking failed: {e}")
        return []


async def embed_chunks_async(chunks: List[str], batch_size: int = 32) -> List[List[float]]:
    """Generate embeddings using E5 Embedding Service."""
    if not chunks:
        return []

    try:
        result = await embedding_client.embed(
            texts=chunks,
            normalize=True,
            batch_size=batch_size,
        )
        embeddings = result.get("embeddings", [])
        dimension = result.get("dimension", 0)
        logger.info(f"Embedded {len(embeddings)} chunks (dimension={dimension})")
        return embeddings

    except Exception as e:
        logger.error(f"Embedding failed: {e}")
        return []


@app.task(name="tasks.document_processing.process_document", bind=True)
def process_document(
    self: Task,
    file_path: str,
    file_hash: str,
    dept_id: int,
    role_id: int
):
    """
    Process a single document using microservices:
    1. Extract text (GLM-OCR for images, libraries for docs)
    2. Chunk text (Hybrid Chunking Service)
    3. Generate embeddings (E5 Embedding Service)
    4. Upsert to Qdrant
    5. Log to PostgreSQL (document + doc_chunk)
    """
    logger.info(f"Processing document: {file_path}")

    loop = asyncio.new_event_loop()
    asyncio.set_event_loop(loop)

    try:
        path = Path(file_path)
        filename = path.name
        file_type = path.suffix.lower()
        file_size = path.stat().st_size

        # Check for existing document
        existing_doc = get_existing_document(file_hash)
        existing_doc_id = existing_doc["doc_id"] if existing_doc else None

        # Upsert document metadata (status = processing)
        doc_id = upsert_document_metadata(
            filename=filename,
            file_path=file_path,
            file_type=file_type,
            file_size=file_size,
            file_hash=file_hash,
            dept_id=dept_id,
            role_id=role_id,
            existing_doc_id=existing_doc_id,
        )

        # Step 1: Extract text
        extracted_text = loop.run_until_complete(extract_text_async(file_path, file_type))
        if not extracted_text:
            logger.warning(f"No text extracted from {file_path}")
            mark_document_status(doc_id, "failed")
            return {"status": "skipped", "reason": "no_text"}

        # Step 2: Chunk text
        chunks = loop.run_until_complete(
            chunk_text_async(extracted_text, method="hybrid", chunk_size=1000, overlap=200)
        )
        if not chunks:
            logger.warning(f"No chunks created from {file_path}")
            mark_document_status(doc_id, "failed")
            return {"status": "skipped", "reason": "no_chunks"}

        # Step 3: Generate embeddings
        embeddings = loop.run_until_complete(embed_chunks_async(chunks, batch_size=32))
        if not embeddings:
            logger.warning(f"No embeddings generated for {file_path}")
            mark_document_status(doc_id, "failed")
            return {"status": "skipped", "reason": "no_embeddings"}
        if len(embeddings) != len(chunks):
            raise RuntimeError(
                f"Chunk/embedding mismatch for {filename}: {len(chunks)} chunks, {len(embeddings)} embeddings"
            )

        # Step 4: Remove old Qdrant points if re-indexing
        qdrant_client = get_qdrant_client()
        collection_name = ensure_qdrant_collection(qdrant_client)

        if existing_doc_id:
            old_point_ids = get_existing_qdrant_ids(existing_doc_id)
            delete_existing_qdrant_points(qdrant_client, collection_name, old_point_ids)
            delete_existing_chunks(existing_doc_id)

        # Step 5: Upsert to Qdrant
        point_ids = upsert_to_qdrant(
            client=qdrant_client,
            collection_name=collection_name,
            doc_id=doc_id,
            chunks=chunks,
            embeddings=embeddings,
            dept_id=dept_id,
            role_id=role_id,
            filename=filename,
            file_path=file_path,
            file_type=file_type,
        )

        # Step 6: Persist chunk metadata in PostgreSQL
        insert_doc_chunks(
            doc_id=doc_id,
            chunks=chunks,
            point_ids=point_ids,
        )

        # Mark document as indexed
        mark_document_status(doc_id, "indexed")

        logger.info(
            f"Successfully processed {filename}: chunks={len(chunks)}, vectors={len(point_ids)}"
        )
        return {
            "status": "completed",
            "filename": filename,
            "doc_id": doc_id,
            "chunks": len(chunks),
            "embeddings": len(embeddings),
            "dept_id": dept_id,
            "role_id": role_id,
        }
    except Exception as e:
        logger.error(f"Failed to process document {file_path}: {e}")
        return {"status": "failed", "error": str(e)}
    finally:
        loop.close()
